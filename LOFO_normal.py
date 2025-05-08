import io
from pprint import pprint
import pandas as pd
import numpy as np
from sklearn.model_selection import train_test_split, GridSearchCV, cross_val_score, KFold
from sklearn.ensemble import RandomForestRegressor, ExtraTreesRegressor, StackingRegressor
from sklearn.linear_model import Ridge, RidgeCV, ElasticNet
from xgboost import XGBRegressor
from lightgbm import LGBMRegressor
from sklearn.neural_network import MLPRegressor
from mlxtend.regressor import StackingCVRegressor
from sklearn.preprocessing import StandardScaler
from sklearn.impute import SimpleImputer
from sklearn.pipeline import Pipeline
from sklearn.metrics import mean_squared_error, r2_score
import networkx as nx
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import pearsonr
from sklearn.impute import KNNImputer
from sklearn.preprocessing import StandardScaler, RobustScaler, FunctionTransformer, MinMaxScaler
from sklearn.compose import ColumnTransformer
from tabulate import tabulate

# If you need python-pptx, you can install it in your local environment
# !pip install python-pptx

#############################
# Step 1: Load the Datasets #
#############################

# Replace 'ecological_data.xlsx' with the path to your ecological dataset
df_ecological = pd.read_excel('Ecological_characterization_Hanna_Markle.xlsx')
print("\n--- First Few Rows of the Ecological Dataset ---\n")
print(df_ecological.head())

# Exploratory Data Analysis function
def explore_data(data):
    """Perform exploratory data analysis on the dataset."""
    print("\n--- Dataset Overview ---\n")
    print(data.info())
    print("\n--- Summary Statistics ---\n")
    print(data.describe())

# Perform EDA
explore_data(df_ecological)

# Read the OTUs associated with pathogens
df_gyrb = pd.read_excel('gyrb_fragment_dataset1.xlsx')  # Replace with correct filename/path
print("\n--- First Few Rows of the Pathogen-Associated OTUs ---\n")
print(df_gyrb.head())

# Read the relative abundance info
df_otu_abundance = pd.read_excel('Leaf_and_Root_OTU_Relative_Abundance.xlsx')  # Replace with correct filename/path
print("\n--- First Few Rows of the OTU Abundance Data ---\n")
print(df_otu_abundance.head())

#################################
# Additional Data Processing    #
#################################
# You can now continue with any preprocessing, merging, modeling, etc.

# Example: Print shapes to verify successful loading
print("\nDataset shapes:")
print("Ecological data shape:", df_ecological.shape)
# print("Pathogen OTUs data shape:", df_gyrb.shape)
print("OTU Abundance data shape:", df_otu_abundance.shape)

"""
CONDITIONAL INDEPENDENCE TESTING USING HARDCODED XGBOOST PARAMETERS
with RECURSIVE FEATURE SELECTION that avoids data leakage by using a holdout set
and handles missing data + robust scaling.
"""

import os
import pandas as pd
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt

from math import sqrt
from sklearn.metrics import r2_score, mean_squared_error
from sklearn.model_selection import train_test_split, KFold
from sklearn.impute import KNNImputer
from sklearn.preprocessing import RobustScaler
import xgboost as xgb

# ---------------------------------------------------------------------
# 1) HELPER: Preprocess and Split (Impute/Drop + Scaling)
# ---------------------------------------------------------------------

tasks = {
    "Microbiota Richness (Root) with Plant OTUs": {
        "target": "richness_microbiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Microbiota Richness (Leaf) with Plant OTUs": {
        "target": "richness_microbiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Microbiota Richness (Root) with Plant Metrics": {
        "target": "richness_microbiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Microbiota Richness (Leaf) with Plant Metrics": {
        "target": "richness_microbiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Microbiota Shannon (Root) with Plant OTUs": {
        "target": "Shannon_microbiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Microbiota Shannon (Leaf) with Plant OTUs": {
        "target": "Shannon_microbiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Microbiota Shannon (Root) with Plant Metrics": {
        "target": "Shannon_microbiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Microbiota Shannon (Leaf) with Plant Metrics": {
        "target": "Shannon_microbiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Pathobiota Richness (Root) with Plant OTUs": {
        "target": "richness_pathobiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Pathobiota Richness (Leaf) with Plant OTUs": {
        "target": "richness_pathobiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Pathobiota Richness (Root) with Plant Metrics": {
        "target": "richness_pathobiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Pathobiota Richness (Leaf) with Plant Metrics": {
        "target": "richness_pathobiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Pathobiota Shannon (Root) with Plant OTUs": {
        "target": "Shannon_pathobiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Pathobiota Shannon (Leaf) with Plant OTUs": {
        "target": "Shannon_pathobiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Pathobiota Shannon (Root) with Plant Metrics": {
        "target": "Shannon_pathobiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Pathobiota Shannon (Leaf) with Plant Metrics": {
        "target": "Shannon_pathobiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    }
}

def preprocess_data_for_task(full_df, features, target, strategy='impute'):
    """
    Cleans and splits the data for a given task:
      1) Remove rows where target is NaN.
      2) Split into train (80%) and holdout (20%).
      3) On the training set:
         - If strategy='drop', drop rows with any missing feature.
           Then re-split if needed (although here we do it after dropping).
         - If strategy='impute', do KNNImputer.
      4) Robust-scale both train and holdout features using scaler fit on train.
      5) If we used imputation, apply the same imputer to the holdout as well.

    Returns
    -------
    X_train : DataFrame
    y_train : Series
    X_holdout : DataFrame
    y_holdout : Series

    The same shape as your original pipeline, but guaranteed no NaNs in the label,
    and scaled/imputed features.
    """
    # 1) Drop rows where target is NaN
    valid_idx = full_df[target].notna()
    df_clean = full_df[valid_idx].copy()  # keep only rows with a valid target

    # 2) Train/Holdout split
    train_df, holdout_df = train_test_split(df_clean, test_size=0.2, random_state=42)
    train_df = train_df.reset_index(drop=True)
    holdout_df = holdout_df.reset_index(drop=True)

    # Extract X, y from the train split
    X_train = train_df[features].copy()
    y_train = train_df[target].copy()

    # Extract X, y from the holdout split
    X_holdout = holdout_df[features].copy()
    y_holdout = holdout_df[target].copy()

    # -- 3) Handle missing values in features --
    if strategy == 'drop':
        # Drop rows in the *train* set that have missing features
        # (We do NOT re-split, we simply lose some train rows)
        train_valid = X_train.notna().all(axis=1)  # row-wise check
        X_train = X_train[train_valid].copy()
        y_train = y_train[train_valid].copy()

        # For holdout, we also drop any row with missing features
        holdout_valid = X_holdout.notna().all(axis=1)
        X_holdout = X_holdout[holdout_valid].copy()
        y_holdout = y_holdout[holdout_valid].copy()

    elif strategy == 'impute':
        # Fit a KNN imputer on training
        imputer = KNNImputer(n_neighbors=5)
        X_train_array = imputer.fit_transform(X_train)
        # Transform holdout with same imputer
        X_holdout_array = imputer.transform(X_holdout)

        # Rebuild DataFrames
        X_train = pd.DataFrame(X_train_array, columns=X_train.columns, index=X_train.index)
        X_holdout = pd.DataFrame(X_holdout_array, columns=X_holdout.columns, index=X_holdout.index)

    else:
        raise ValueError("Invalid strategy. Must be 'impute' or 'drop'.")

    # -- 4) Robust scaling both sets based on train statistics
    scaler = RobustScaler()
    X_train_scaled = pd.DataFrame(
        scaler.fit_transform(X_train),
        columns=X_train.columns,
        index=X_train.index
    )
    X_holdout_scaled = pd.DataFrame(
        scaler.transform(X_holdout),
        columns=X_holdout.columns,
        index=X_holdout.index
    )

    return X_train_scaled, y_train, X_holdout_scaled, y_holdout


# ---------------------------------------------------------------------
# 2) XGBoost Training with Early Stopping (Hardcoded Params)
# ---------------------------------------------------------------------
def xgb_train_early_stopping(X, y, params, early_stopping_rounds=20, num_rounds=500, random_state=42):
    """
    Trains an XGBoost model with early stopping using the provided hardcoded parameters.
    This function performs its own split into training and validation sets (for early stopping)
    on the provided data (assumed to be the training set from the outer split).
    
    Returns
    -------
    model : xgboost.Booster
    rmse_train : float
    rmse_val   : float
    """
    # Reset index to ensure alignment
    X = X.reset_index(drop=True)
    y = y.reset_index(drop=True)
    
    # Split into training and validation for early stopping
    X_train, X_val, y_train, y_val = train_test_split(
        X, y, test_size=0.2, random_state=random_state
    )
    
    dtrain = xgb.DMatrix(X_train, label=y_train, feature_names=X_train.columns.tolist())
    dval   = xgb.DMatrix(X_val,   label=y_val,   feature_names=X_val.columns.tolist())
    
    watchlist = [(dtrain, 'train'), (dval, 'eval')]
    model = xgb.train(
        params=params,
        dtrain=dtrain,
        num_boost_round=num_rounds,
        evals=watchlist,  # pass evals explicitly
        early_stopping_rounds=early_stopping_rounds,
        verbose_eval=False
    )

    # Predictions on training and validation sets
    y_pred_train = model.predict(xgb.DMatrix(X_train, feature_names=X_train.columns.tolist()))
    y_pred_val   = model.predict(xgb.DMatrix(X_val,   feature_names=X_val.columns.tolist()))
    
    rmse_train = sqrt(mean_squared_error(y_train, y_pred_train))
    rmse_val   = sqrt(mean_squared_error(y_val,   y_pred_val))
    
    return model, rmse_train, rmse_val

# ---------------------------------------------------------------------
# 3) Nested CV: compute_cv_rmse
# ---------------------------------------------------------------------
def compute_cv_rmse(X, y, params, features, n_splits=5):
    """
    Helper to compute mean CV RMSE with the given subset of features on X, y.
    Each fold calls xgb_train_early_stopping to get a validation RMSE.
    """
    kf = KFold(n_splits=n_splits, shuffle=True, random_state=42)
    rmse_list = []
    for train_idx, val_idx in kf.split(X):
        X_tr = X.iloc[train_idx][features]
        y_tr = y.iloc[train_idx]
        # We use the existing early-stopping function to fit on each fold
        _, _, rmse_val = xgb_train_early_stopping(X_tr, y_tr, params)
        rmse_list.append(rmse_val)
    return np.mean(rmse_list)

# ---------------------------------------------------------------------
# 4) Recursive Feature Selection (Nested CV on Training Set Only)
# ---------------------------------------------------------------------
def recursive_feature_selection(
    X_train, y_train, params,
    current_features=None,
    drop_history=None,
    n_splits=5,
    rmse_history=None
):
    """
    Recursively drops one feature at a time (using nested CV on the training set only)
    to determine if dropping a feature improves CV RMSE.
    
    Returns
    -------
    best_features : list
        The best subset of features found.
    rmse_history : list of dict
        Each dict can store: {'features': [...], 'rmse': <float>}
    """
    if current_features is None:
        current_features = X_train.columns.tolist()
    if drop_history is None:
        drop_history = []
    if rmse_history is None:
        rmse_history = []

    # 1) Compute baseline CV RMSE
    baseline_rmse = compute_cv_rmse(X_train, y_train, params, current_features, n_splits)
    print(f"[RecursiveFS] Current features ({len(current_features)}): {current_features}")
    print(f"[RecursiveFS] Baseline CV RMSE = {baseline_rmse:.3f}")

    # Record in rmse_history
    rmse_history.append({
        'features': current_features.copy(),
        'rmse': baseline_rmse
    })

    best_features = current_features
    best_rmse = baseline_rmse
    improved = False
    dropped_feature = None

    # 2) Try dropping each feature
    for feat in current_features:
        candidate = [f for f in current_features if f != feat]
        if not candidate:
            continue

        candidate_rmse = compute_cv_rmse(X_train, y_train, params, candidate, n_splits)
        print(f"Candidate drop {feat}: CV RMSE = {candidate_rmse:.3f} (Baseline = {baseline_rmse:.3f})")

        if candidate_rmse < (best_rmse - 1e-6):
            best_features = candidate
            best_rmse = candidate_rmse
            improved = True
            dropped_feature = feat

    # 3) If improved, drop that feature and recurse
    if improved:
        print(f"Improvement found: dropping '{dropped_feature}' => new CV RMSE = {best_rmse:.3f}")
        drop_history.append(dropped_feature)
        return recursive_feature_selection(
            X_train, y_train, params,
            current_features=best_features,
            drop_history=drop_history,
            n_splits=n_splits,
            rmse_history=rmse_history
        )
    else:
        print("No further improvement found.")
        return best_features, rmse_history

# ---------------------------------------------------------------------
# 5) Final Model Evaluation on the Holdout Test Set
# ---------------------------------------------------------------------
def final_model_evaluation(
    X_train, y_train,
    X_holdout, y_holdout,
    params,
    early_stopping_rounds=25,
    num_rounds=500
):
    """
    Trains a final model on the entire training set (using the selected features)
    and evaluates it on the holdout test set.
    
    Returns
    -------
    model : xgboost.Booster
    rmse_train : float
    rmse_holdout : float
    r2_holdout : float
    test_std : float
    """
    model, rmse_train, _ = xgb_train_early_stopping(
        X_train, y_train,
        params,
        early_stopping_rounds=early_stopping_rounds,
        num_rounds=num_rounds
    )
    d_holdout = xgb.DMatrix(X_holdout, feature_names=X_holdout.columns.tolist())
    y_pred_holdout = model.predict(d_holdout)
    rmse_holdout = sqrt(mean_squared_error(y_holdout, y_pred_holdout))
    r2_holdout   = r2_score(y_holdout, y_pred_holdout)
    test_std     = np.std(y_holdout)

    return model, rmse_train, rmse_holdout, r2_holdout, test_std

# ---------------------------------------------------------------------
# 6) MAIN: Perform RFS + Evaluate
# ---------------------------------------------------------------------
def perform_recursive_feature_selection(framework, tasks, strategy='impute'):
    """
    For each 'Richness' task in tasks:
      1) Preprocess data (drop/impute + scale) => X_train_full, y_train_full, X_holdout, y_holdout
      2) Recursively (via nested CV) drop features if it improves CV RMSE
      3) Train final model on best subset + measure holdout performance

    Parameters
    ----------
    framework : object
        Must have .data (DataFrame).
    tasks : dict
        {task_name: {'target':..., 'features':...}, ...}
    strategy : str
        'impute' => KNNImputer on missing features,
        'drop'   => drop rows with any missing feature.

    Returns
    -------
    results : dict
        Maps task_name => {
          'Selected_Features': ...,
          'CV_RMSE_History': [...],
          'Final_Train_RMSE': ...,
          'Final_Test_RMSE': ...,
          'Final_R2': ...,
          'Final_Test_Std': ...
        }
    """
    results = {}

    # Filter tasks to those with "Richness" in name or target
    richness_tasks = {
        t_name: t_info
        for t_name, t_info in tasks.items()
        if "Richness" in t_name or "richness" in t_info['target']
    }

    for task_name, task_info in richness_tasks.items():
        print(f"\n=== Processing Task: {task_name} ===")
        target   = task_info['target']
        features = task_info['features']

        # ---------------------------
        # 1) Preprocess data
        # ---------------------------
        # Merge drop/impute + scaling + train/holdout splitting
        X_train_full, y_train_full, X_holdout, y_holdout = preprocess_data_for_task(
            framework.data, features, target, strategy=strategy
        )

        if len(X_train_full) == 0 or len(X_holdout) == 0:
            print(f"[WARNING] No valid data left for {task_name} after '{strategy}'. Skipping...")
            continue

        # Hardcoded XGBoost params for each task
        hardcoded_params = {
            "Microbiota Richness (Root) with Plant OTUs": {
                'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5,
                'min_child_weight': 1, 'reg_alpha': 0.5, 'reg_lambda': 0.5,
                'subsample': 0.8, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Microbiota Richness (Leaf) with Plant OTUs": {
                'colsample_bytree': 0.8, 'learning_rate': 0.01, 'max_depth': 5,
                'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 2.0,
                'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Microbiota Richness (Root) with Plant Metrics": {
                'colsample_bytree': 0.9, 'learning_rate': 0.1, 'max_depth': 3,
                'min_child_weight': 1, 'reg_alpha': 0.1, 'reg_lambda': 1.0,
                'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Microbiota Richness (Leaf) with Plant Metrics": {
                'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 4,
                'min_child_weight': 1, 'reg_alpha': 0.1, 'reg_lambda': 1.0,
                'subsample': 0.8, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Microbiota Shannon (Root) with Plant OTUs": {
                'colsample_bytree': 0.8, 'learning_rate': 0.01, 'max_depth': 5,
                'min_child_weight': 5, 'reg_alpha': 1.0, 'reg_lambda': 0.5,
                'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Microbiota Shannon (Leaf) with Plant OTUs": {
                'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5,
                'min_child_weight': 3, 'reg_alpha': 0.1, 'reg_lambda': 0.5,
                'subsample': 0.8, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Microbiota Shannon (Root) with Plant Metrics": {
                'colsample_bytree': 0.9, 'learning_rate': 0.01, 'max_depth': 4,
                'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 0.5,
                'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Microbiota Shannon (Leaf) with Plant Metrics": {
                'colsample_bytree': 0.7, 'learning_rate': 0.05, 'max_depth': 3,
                'min_child_weight': 3, 'reg_alpha': 0.1, 'reg_lambda': 0.5,
                'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Pathobiota Richness (Root) with Plant OTUs": {
                'colsample_bytree': 0.9, 'learning_rate': 0.01, 'max_depth': 5,
                'min_child_weight': 1, 'reg_alpha': 0.5, 'reg_lambda': 2.0,
                'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Pathobiota Richness (Leaf) with Plant OTUs": {
                'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5,
                'min_child_weight': 1, 'reg_alpha': 0.5, 'reg_lambda': 2.0,
                'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Pathobiota Richness (Root) with Plant Metrics": {
                'colsample_bytree': 0.9, 'learning_rate': 0.05, 'max_depth': 5,
                'min_child_weight': 5, 'reg_alpha': 1.0, 'reg_lambda': 2.0,
                'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Pathobiota Richness (Leaf) with Plant Metrics": {
                'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5,
                'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 2.0,
                'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Pathobiota Shannon (Root) with Plant OTUs": {
                'colsample_bytree': 0.7, 'learning_rate': 0.1, 'max_depth': 4,
                'min_child_weight': 5, 'reg_alpha': 1.0, 'reg_lambda': 1.0,
                'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Pathobiota Shannon (Leaf) with Plant OTUs": {
                'colsample_bytree': 0.9, 'learning_rate': 0.01, 'max_depth': 3,
                'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 2.0,
                'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Pathobiota Shannon (Root) with Plant Metrics": {
                'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5,
                'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 2.0,
                'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
            },
            "Pathobiota Shannon (Leaf) with Plant Metrics": {
                'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 4,
                'min_child_weight': 5, 'reg_alpha': 0.1, 'reg_lambda': 2.0,
                'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
            }
        }
        params = hardcoded_params.get(task_name, {
            'objective': 'reg:squarederror',
            'max_depth': 3,
            'learning_rate': 0.05,
            'min_child_weight': 5,
            'subsample': 0.8,
            'colsample_bytree': 0.8,
            'gamma': 1.0,
            'reg_alpha': 0.5,
            'reg_lambda': 1.0,
            'seed': 42
        })

        # ---------------------------
        # 2) Recursive Feature Selection (Nested CV)
        # ---------------------------
        best_features, cv_history = recursive_feature_selection(
            X_train_full, y_train_full, params
        )
        print(f"Task: {task_name} -> Best Subset: {best_features}")
        print(f"CV RMSE History: {cv_history}")

        # ---------------------------
        # 3) Final Model on holdout
        # ---------------------------
        final_model, final_train_rmse, final_test_rmse, final_r2, final_test_std = final_model_evaluation(
            X_train_full[best_features], y_train_full,
            X_holdout[best_features],  y_holdout,
            params
        )
        print(f"Final Model Performance on Holdout Test Set for {task_name}:")
        print(f"  Train RMSE: {final_train_rmse:.3f}")
        print(f"  Test RMSE : {final_test_rmse:.3f}")
        print(f"  R2        : {final_r2:.3f}")
        print(f"  Test Std  : {final_test_std:.3f}")

        results[task_name] = {
            "Selected_Features": best_features,
            "CV_RMSE_History": cv_history,
            "Final_Train_RMSE": final_train_rmse,
            "Final_Test_RMSE": final_test_rmse,
            "Final_R2": final_r2,
            "Final_Test_Std": final_test_std
        }

    return results


# ---------------------------------------------------------------------
# POWERPOINT REPORT GENERATION (including final selected features)
# ---------------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_powerpoint_report(results, output_pptx='report.pptx'):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Ecodynamics AI - Final Model Results"
    slide.placeholders[1].text = "Recursive Feature Selection with XGBoost"
    
    bullet_slide_layout = prs.slide_layouts[1]
    
    for task_name, res in results.items():
        slide_task = prs.slides.add_slide(bullet_slide_layout)
        slide_task.shapes.title.text = f"Task: {task_name}"
        tf = slide_task.shapes.placeholders[1].text_frame
        tf.text = (
            f"Final Test RMSE: {res['Final_Test_RMSE']:.3f}\n"
            f"Final R2: {res['Final_R2']:.3f}\n"
            f"Test Std: {res['Final_Test_Std']:.3f}\n"
            f"Selected Features ({len(res['Selected_Features'])}): {', '.join(res['Selected_Features'])}"
        )
    
    prs.save(output_pptx)
    print(f"PowerPoint saved to {output_pptx}")


# ---------------------------------------------------------------------
# USAGE EXAMPLE (in your real code):
# ---------------------------------------------------------------------
# Suppose:
#   framework.data is your main DataFrame
#   tasks is a dictionary of tasks, e.g.:
#     tasks = {
#       "Microbiota Richness (Root) with Plant OTUs": {
#           "target": "some_target_col",
#           "features": [...]
#       },
#       ...
#     }
#
# Then simply call:
#


from ecodynamics_multimodel_framework import EcodynamicsAI

framework = EcodynamicsAI()
framework.load_data()
framework.load_otu_data()
framework.categorize_features()
framework.filter_columns_of_interest()
framework.validate_and_assign_otu_data()

results = perform_recursive_feature_selection(framework, tasks, strategy='drop')