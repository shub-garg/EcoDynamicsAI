import io
from pprint import pprint
import pandas as pd
import numpy as np
from sklearn.model_selection import train_test_split, GridSearchCV, cross_val_score, KFold
from sklearn.ensemble import RandomForestRegressor, ExtraTreesRegressor, StackingRegressor
from sklearn.linear_model import Ridge, RidgeCV, ElasticNet
from xgboost import XGBRegressor
from lightgbm import LGBMRegressor
from sklearn.neural_network import MLPRegressor
from mlxtend.regressor import StackingCVRegressor
from sklearn.preprocessing import StandardScaler
from sklearn.impute import SimpleImputer
from sklearn.pipeline import Pipeline
from sklearn.metrics import mean_squared_error, r2_score
import networkx as nx
import matplotlib.pyplot as plt
import seaborn as sns
from scipy.stats import pearsonr
from sklearn.impute import KNNImputer
from sklearn.preprocessing import StandardScaler, RobustScaler, FunctionTransformer, MinMaxScaler
from sklearn.compose import ColumnTransformer
from tabulate import tabulate

# If you need python-pptx, you can install it in your local environment
# !pip install python-pptx

#############################
# Step 1: Load the Datasets #
#############################

# Replace 'ecological_data.xlsx' with the path to your ecological dataset
df_ecological = pd.read_excel('Ecological_characterization_Hanna_Markle.xlsx')
print("\n--- First Few Rows of the Ecological Dataset ---\n")
print(df_ecological.head())

# Exploratory Data Analysis function
def explore_data(data):
    """Perform exploratory data analysis on the dataset."""
    print("\n--- Dataset Overview ---\n")
    print(data.info())
    print("\n--- Summary Statistics ---\n")
    print(data.describe())

# Perform EDA
explore_data(df_ecological)

# Read the OTUs associated with pathogens
df_gyrb = pd.read_excel('gyrb_fragment_dataset1.xlsx')  # Replace with correct filename/path
print("\n--- First Few Rows of the Pathogen-Associated OTUs ---\n")
print(df_gyrb.head())

# Read the relative abundance info
df_otu_abundance = pd.read_excel('Leaf_and_Root_OTU_Relative_Abundance.xlsx')  # Replace with correct filename/path
print("\n--- First Few Rows of the OTU Abundance Data ---\n")
print(df_otu_abundance.head())

#################################
# Additional Data Processing    #
#################################
# You can now continue with any preprocessing, merging, modeling, etc.

# Example: Print shapes to verify successful loading
print("\nDataset shapes:")
print("Ecological data shape:", df_ecological.shape)
# print("Pathogen OTUs data shape:", df_gyrb.shape)
print("OTU Abundance data shape:", df_otu_abundance.shape)

"""
CONDITIONAL INDEPENDENCE TESTING USING HARDCODED XGBOOST PARAMETERS
"""

import os
import pandas as pd
import numpy as np
import seaborn as sns
import matplotlib.pyplot as plt
from sklearn.metrics import r2_score, mean_squared_error
from math import sqrt
from sklearn.model_selection import train_test_split, KFold
from scipy.stats import ttest_rel
import xgboost as xgb

# ---------------------------------------------------------------------
# EXTERNAL FUNCTION: XGBOOST TRAINING WITH EARLY STOPPING USING HARD-CODED PARAMETERS
# ---------------------------------------------------------------------
def xgb_train_early_stopping(X, y, params, early_stopping_rounds=20, num_rounds=500, random_state=42):
    """
    Trains an XGBoost model with early stopping using the provided hardcoded parameters.
    
    Parameters:
      X : pandas.DataFrame
          Feature matrix.
      y : pandas.Series
          Target variable.
      params : dict
          Hardcoded parameters for XGBoost.
      early_stopping_rounds : int, optional
          Early stopping rounds.
      num_rounds : int, optional
          Maximum number of boosting rounds.
      random_state : int, optional
          Random state for reproducibility.
    
    Returns:
      model : trained XGBoost Booster
      X_test : Test set features (DataFrame)
      y_test : Test set targets (Series)
      rmse_train : RMSE on the training set
      rmse_test : RMSE on the test set
    """
    # Reset index to ensure alignment
    X = X.reset_index(drop=True)
    y = y.reset_index(drop=True)
    
    # Split data into train and test sets
    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=random_state)
    X_train = X_train.reset_index(drop=True)
    X_test = X_test.reset_index(drop=True)
    y_train = y_train.reset_index(drop=True)
    y_test = y_test.reset_index(drop=True)
    
    # Further split training into sub-training and validation sets for early stopping
    X_train_sub, X_val, y_train_sub, y_val = train_test_split(X_train, y_train, test_size=0.2, random_state=random_state)
    
    # Create DMatrices with feature names (as list) so that Booster returns proper keys
    dtrain = xgb.DMatrix(X_train_sub, label=y_train_sub, feature_names=X_train_sub.columns.tolist())
    dval = xgb.DMatrix(X_val, label=y_val, feature_names=X_val.columns.tolist())
    dtest = xgb.DMatrix(X_test, feature_names=X_test.columns.tolist())
    
    watchlist = [(dtrain, 'train'), (dval, 'eval')]
    
    print("Starting training with early stopping using hardcoded parameters...")
    model = xgb.train(params, dtrain, num_rounds, watchlist,
                      early_stopping_rounds=early_stopping_rounds, verbose_eval=False)
    
    # Predict on the entire training set and test set
    y_pred_train = model.predict(xgb.DMatrix(X_train, feature_names=X_train.columns.tolist()))
    y_pred_test = model.predict(dtest)
    
    rmse_train = sqrt(mean_squared_error(y_train, y_pred_train))
    rmse_test = sqrt(mean_squared_error(y_test, y_pred_test))
    
    return model, X_test, y_test, rmse_train, rmse_test

# ---------------------------------------------------------------------
# LOFO and Correlation Plotting Function (using XGBoost baseline)
# ---------------------------------------------------------------------
def perform_lofo_and_correlation_plots(framework, tasks, top_k_corr=15,
                                       do_significance_test=True, test_mode=False,
                                       test_tasks_limit=1, figs_dir='figs'):
    """
    For each 'Richness' task, uses XGBoost with early stopping (with hardcoded parameters)
    as the baseline model, plots correlation heatmaps, performs LOFO analysis on all features,
    and returns dictionaries with results and figure paths.
    """
    
    if test_mode:
        subset_task_names = list(tasks.keys())[:test_tasks_limit]
        tasks = {k: tasks[k] for k in subset_task_names}
        print(f"[TEST MODE] Only running LOFO on the first {test_tasks_limit} tasks.")
    
    # Filter tasks to those involving 'Richness'
    richness_tasks = { t_name: t_info for t_name, t_info in tasks.items()
                       if "Richness" in t_name or "richness" in t_info['target'] }
    
    lofo_results = {}
    correlation_figs = {}
    feature_importance_figs = {}
    dataset_sizes = {}
    
    # Define hardcoded parameters for each task
    hardcoded_params = {
        "Microbiota Richness (Root) with Plant OTUs": {
            'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5, 
            'min_child_weight': 1, 'reg_alpha': 0.5, 'reg_lambda': 0.5, 
            'subsample': 0.8, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Microbiota Richness (Leaf) with Plant OTUs": {
            'colsample_bytree': 0.8, 'learning_rate': 0.01, 'max_depth': 5, 
            'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 2.0, 
            'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Microbiota Richness (Root) with Plant Metrics": {
            'colsample_bytree': 0.9, 'learning_rate': 0.1, 'max_depth': 3, 
            'min_child_weight': 1, 'reg_alpha': 0.1, 'reg_lambda': 1.0, 
            'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Microbiota Richness (Leaf) with Plant Metrics": {
            'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 4, 
            'min_child_weight': 1, 'reg_alpha': 0.1, 'reg_lambda': 1.0, 
            'subsample': 0.8, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Microbiota Shannon (Root) with Plant OTUs": {
            'colsample_bytree': 0.8, 'learning_rate': 0.01, 'max_depth': 5, 
            'min_child_weight': 5, 'reg_alpha': 1.0, 'reg_lambda': 0.5, 
            'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Microbiota Shannon (Leaf) with Plant OTUs": {
            'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5, 
            'min_child_weight': 3, 'reg_alpha': 0.1, 'reg_lambda': 0.5, 
            'subsample': 0.8, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Microbiota Shannon (Root) with Plant Metrics": {
            'colsample_bytree': 0.9, 'learning_rate': 0.01, 'max_depth': 4, 
            'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 0.5, 
            'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Microbiota Shannon (Leaf) with Plant Metrics": {
            'colsample_bytree': 0.7, 'learning_rate': 0.05, 'max_depth': 3, 
            'min_child_weight': 3, 'reg_alpha': 0.1, 'reg_lambda': 0.5, 
            'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Pathobiota Richness (Root) with Plant OTUs": {
            'colsample_bytree': 0.9, 'learning_rate': 0.01, 'max_depth': 5, 
            'min_child_weight': 1, 'reg_alpha': 0.5, 'reg_lambda': 2.0, 
            'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Pathobiota Richness (Leaf) with Plant OTUs": {
            'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5, 
            'min_child_weight': 1, 'reg_alpha': 0.5, 'reg_lambda': 2.0, 
            'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Pathobiota Richness (Root) with Plant Metrics": {
            'colsample_bytree': 0.9, 'learning_rate': 0.05, 'max_depth': 5, 
            'min_child_weight': 5, 'reg_alpha': 1.0, 'reg_lambda': 2.0, 
            'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Pathobiota Richness (Leaf) with Plant Metrics": {
            'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5, 
            'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 2.0, 
            'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Pathobiota Shannon (Root) with Plant OTUs": {
            'colsample_bytree': 0.7, 'learning_rate': 0.1, 'max_depth': 4, 
            'min_child_weight': 5, 'reg_alpha': 1.0, 'reg_lambda': 1.0, 
            'subsample': 0.9, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Pathobiota Shannon (Leaf) with Plant OTUs": {
            'colsample_bytree': 0.9, 'learning_rate': 0.01, 'max_depth': 3, 
            'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 2.0, 
            'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Pathobiota Shannon (Root) with Plant Metrics": {
            'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 5, 
            'min_child_weight': 1, 'reg_alpha': 1.0, 'reg_lambda': 2.0, 
            'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
        },
        "Pathobiota Shannon (Leaf) with Plant Metrics": {
            'colsample_bytree': 0.7, 'learning_rate': 0.01, 'max_depth': 4, 
            'min_child_weight': 5, 'reg_alpha': 0.1, 'reg_lambda': 2.0, 
            'subsample': 0.7, 'objective': 'reg:squarederror', 'seed': 42
        }
    }
    
    for task_name, task_info in richness_tasks.items():
        print(f"\n=== LOFO & Correlation Plot for {task_name} ===")
        target = task_info['target']
        all_features = task_info['features']
    
        # Prepare data from framework (assumes framework.data is a DataFrame)
        X_full = framework.data[all_features]
        y_full = framework.data[target]
    
        # Use the framework's imputation (and scaling) method
        X_train, X_test, y_train, y_test = framework.impute_or_drop(X_full, y_full)
        X_train = X_train.reset_index(drop=True)
        X_test = X_test.reset_index(drop=True)
        y_train = y_train.reset_index(drop=True)
        y_test = y_test.reset_index(drop=True)
        X_train = X_train.apply(pd.to_numeric, errors='coerce')
        X_test = X_test.apply(pd.to_numeric, errors='coerce')
    
        num_rows, num_cols = X_train.shape
        dataset_sizes[task_name] = (num_rows, num_cols + 1)
    
        # Retrieve hardcoded parameters for the task; if not found, use default parameters
        params = hardcoded_params.get(task_name, {
            'objective': 'reg:squarederror',
            'max_depth': 3,
            'learning_rate': 0.05,
            'min_child_weight': 5,
            'subsample': 0.8,
            'colsample_bytree': 0.8,
            'gamma': 1.0,
            'reg_alpha': 0.5,
            'reg_lambda': 1.0,
            'seed': 42
        })
    
        # Baseline model using XGBoost with early stopping
        model, X_test_model, y_test_model, rmse_train, rmse_test = xgb_train_early_stopping(
            X_train, y_train, params, early_stopping_rounds=20
        )
        dtest = xgb.DMatrix(X_test, feature_names=X_test.columns.tolist())
        y_pred_test = model.predict(dtest)
        baseline_r2 = r2_score(y_test, y_pred_test)
        baseline_rmse = sqrt(mean_squared_error(y_test, y_pred_test))
        baseline_std = np.std(y_test)
        print(f"Task: {task_name}")
        print(f"Full Feature Model - RMSE Train: {rmse_train:.3f}, RMSE Test: {baseline_rmse:.3f}")
    
        # Get feature importances from the Booster (using gain)
        importance_dict = model.get_score(importance_type='gain')
        importance_df = pd.DataFrame(list(importance_dict.items()), columns=['Feature', 'Importance'])
        importance_df = importance_df.sort_values("Importance", ascending=False)
    
        # Plot and save feature importance
        plt.figure(figsize=(6, 4))
        sns.barplot(x="Importance", y="Feature", data=importance_df.head(15))
        plt.title(f"Importance - {task_name}")
        plt.tight_layout()
        safe_task_name = task_name.replace(" ", "_").replace("(", "").replace(")", "").replace("/", "")
        featimp_path = os.path.join(figs_dir, f"{safe_task_name}_feature_importance.png")
        # plt.savefig(featimp_path, dpi=150)
        plt.close()
        feature_importance_figs[task_name] = featimp_path
    
        # Plot correlation heatmap for the top_k_corr features (as per importance order)
        top_features_for_corr = importance_df.head(top_k_corr)['Feature'].tolist()
        top_features_for_corr = [f for f in top_features_for_corr if f in X_train.columns]
        corr_matrix = X_train[top_features_for_corr].corr()
        plt.figure(figsize=(8,6))
        sns.heatmap(corr_matrix, annot=True, cmap='RdBu', center=0)
        plt.title(f"Correlation (Top {top_k_corr} Features) - {task_name}")
        plt.tight_layout()
        fig_path = os.path.join(figs_dir, f"{safe_task_name}_corr_heatmap.png")
        # plt.savefig(fig_path, dpi=150)
        plt.show()
        plt.close()
        correlation_figs[task_name] = fig_path
    
        # LOFO analysis using XGBoost with early stopping over ALL features
        all_feats = importance_df['Feature'].tolist()  # use all features
        drop_results = []
    
        for feat in all_feats:
            print(f"\nDropping feature '{feat}' => measuring test performance difference...")
            X_train_drop = X_train.drop(columns=[feat], errors='ignore')
            X_test_drop = X_test.drop(columns=[feat], errors='ignore')
    
            model_drop, X_test_drop_model, y_test_drop, rmse_train_drop, rmse_test_drop = xgb_train_early_stopping(
                X_train_drop, y_train, params, early_stopping_rounds=20
            )
            dtest_drop = xgb.DMatrix(X_test_drop, feature_names=X_test_drop.columns.tolist())
            y_pred_drop = model_drop.predict(dtest_drop)
            drop_r2 = r2_score(y_test, y_pred_drop)
            drop_rmse = sqrt(mean_squared_error(y_test, y_pred_drop))
    
            r2_diff = baseline_r2 - drop_r2
            rmse_diff = drop_rmse - baseline_rmse
            print(f"   => Original R2: {baseline_r2:.3f}, Dropped R2: {drop_r2:.3f}, ΔR² = {r2_diff:.3f}")
            print(f"   => Original RMSE: {baseline_rmse:.3f}, Dropped RMSE: {drop_rmse:.3f}, ΔRMSE = {rmse_diff:.3f}")
    
            if do_significance_test:
                sig_result = lofo_significance_test(
                    X_train, y_train, feature_to_drop=feat, n_splits=5, n_repeats=3
                )
                p_val = sig_result['p_val']
                mean_diff_cv = sig_result['mean_diff']
                print(f"   => Repeated-CV: mean R2 diff = {mean_diff_cv:.3f}, p-value = {p_val:.4g}")
            else:
                p_val = None
                mean_diff_cv = None
    
            drop_results.append({
                "feature": feat,
                "Test_R2_after_drop": drop_r2,
                "Test_R2_diff": r2_diff,
                "Test_RMSE_after_drop": drop_rmse,
                "Test_RMSE_diff": rmse_diff,
                "Train_RMSE_after_drop": rmse_train_drop,
                "Test_Std": np.std(y_test),
                "RepeatedCV_mean_diff": mean_diff_cv,
                "RepeatedCV_p_value": p_val
            })
    
        # Determine best-drop(s): find the drop with lowest test RMSE
        if drop_results:
            drop_rmses = [d["Test_RMSE_after_drop"] for d in drop_results]
            min_rmse = min(drop_rmses)
            tolerance = 1e-6  # use a small tolerance
            best_drops = [d for d in drop_results if abs(d["Test_RMSE_after_drop"] - min_rmse) < tolerance]
        else:
            best_drops = []
    
        lofo_result_for_task = {
            "Baseline R2": baseline_r2,
            "Baseline RMSE": baseline_rmse,
            "Baseline Test Std": baseline_std,
            "Feature_Drops": drop_results,
            "Best_Drops": best_drops
        }
    
        lofo_results[task_name] = lofo_result_for_task
    
    return lofo_results, correlation_figs, dataset_sizes, feature_importance_figs

# ---------------------------------------------------------------------
# LOFO Significance Test Function (unchanged)
# ---------------------------------------------------------------------
def lofo_significance_test(X, y, feature_to_drop, n_splits=5, n_repeats=3):
    r2_full_scores = []
    r2_drop_scores = []
    for random_seed in range(n_repeats):
        kf = KFold(n_splits=n_splits, shuffle=True, random_state=42 + random_seed)
        for train_idx, test_idx in kf.split(X):
            X_train_cv, X_test_cv = X.iloc[train_idx], X.iloc[test_idx]
            y_train_cv, y_test_cv = y.iloc[train_idx], y.iloc[test_idx]
            model_full = RandomForestRegressor(
                n_estimators=400, random_state=42, max_features='sqrt',
                max_depth=10, min_samples_split=5, min_samples_leaf=3
            )
            model_full.fit(X_train_cv, y_train_cv)
            y_pred_full = model_full.predict(X_test_cv)
            r2_full = r2_score(y_test_cv, y_pred_full)
            r2_full_scores.append(r2_full)
    
            X_train_dropped = X_train_cv.drop(columns=[feature_to_drop], errors='ignore')
            X_test_dropped = X_test_cv.drop(columns=[feature_to_drop], errors='ignore')
            model_drop = RandomForestRegressor(
                n_estimators=400, random_state=42, max_features='sqrt',
                max_depth=10, min_samples_split=5, min_samples_leaf=3
            )
            model_drop.fit(X_train_dropped, y_train_cv)
            y_pred_drop = model_drop.predict(X_test_dropped)
            r2_drop = r2_score(y_test_cv, y_pred_drop)
            r2_drop_scores.append(r2_drop)
    r2_full_scores = np.array(r2_full_scores)
    r2_drop_scores = np.array(r2_drop_scores)
    r2_diff = r2_full_scores - r2_drop_scores
    mean_diff = np.mean(r2_diff)
    t_stat, p_val = ttest_rel(r2_full_scores, r2_drop_scores)
    return {
        'mean_diff': mean_diff,
        'p_val': p_val,
        'full_model_mean_r2': np.mean(r2_full_scores),
        'drop_model_mean_r2': np.mean(r2_drop_scores),
        'n_splits': n_splits,
        'n_repeats': n_repeats
    }

# ---------------------------------------------------------------------
# PowerPoint Report Generation
# ---------------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_powerpoint_report(lofo_results, correlation_figs=None, dataset_sizes=None,
                             feature_importance_figs=None, output_pptx='report.pptx'):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Ecodynamics AI - Models Results"
    slide.placeholders[1].text = "Data Analysis"
    # Dataset sizes slide
    bullet_slide_layout = prs.slide_layouts[1]
    ds_slide = prs.slides.add_slide(bullet_slide_layout)
    ds_slide.shapes.title.text = "Dataset Sizes"
    ds_body = ds_slide.shapes.placeholders[1].text_frame
    ds_body.text = "Each task uses a dataset with the following shape (X_train):\n"
    for task_name, (rows, cols) in dataset_sizes.items():
        p = ds_body.add_paragraph()
        p.level = 1
        p.text = f"{task_name}: {rows} rows, {cols} columns"
    # Feature importance slides
    image_slide_layout = prs.slide_layouts[5]
    for task_name, fig_path in feature_importance_figs.items():
        if not os.path.exists(fig_path):
            print(f"[WARNING] Feature Importance file not found: {fig_path}")
            continue
        slide_fi = prs.slides.add_slide(image_slide_layout)
        slide_fi.shapes.title.text = f"Feature Importance: {task_name}"
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(4)
        slide_fi.shapes.add_picture(fig_path, left, top, height=height)
    # LOFO results slides
    for task_name, lofo_data in lofo_results.items():
        slide_lofo = prs.slides.add_slide(bullet_slide_layout)
        slide_lofo.shapes.title.text = f"LOFO Results: {task_name}"
        tf_lofo = slide_lofo.shapes.placeholders[1].text_frame
        baseline_rmse = lofo_data["Baseline RMSE"]
        tf_lofo.text = f"Baseline RMSE: {baseline_rmse:.3f}"
        # Report best-drop features
        best_drops = lofo_data.get("Best_Drops", [])
        if best_drops:
            drop_texts = [f"{d['feature']} (Train RMSE: {d['Train_RMSE_after_drop']:.3f}, Test RMSE: {d['Test_RMSE_after_drop']:.3f}, Test Std: {d['Test_Std']:.3f})" for d in best_drops]
            tf_lofo.text += "\nBest Drop(s): " + ", ".join(drop_texts)
        # List all drop details (optional)
        drops = lofo_data["Feature_Drops"]
        for drop_info in drops:
            p = tf_lofo.add_paragraph()
            p.level = 1
            feat = drop_info["feature"]
            rmse_diff = drop_info["Test_RMSE_diff"]
            rmse_after = drop_info["Test_RMSE_after_drop"]
            cv_diff = drop_info.get("RepeatedCV_mean_diff", None)
            p_val = drop_info.get("RepeatedCV_p_value", None)
            text_str = (f"Drop '{feat}' => RMSE diff: {rmse_diff:+.3f} (RMSE after drop: {rmse_after:.3f})")
            if cv_diff is not None and p_val is not None:
                text_str += f" | p = {p_val:.3g}"
            p.text = text_str
    # Correlation heatmap slides
    for task_name, fig_path in correlation_figs.items():
        if not os.path.exists(fig_path):
            print(f"[WARNING] Correlation fig not found: {fig_path}")
            continue
        slide_corr = prs.slides.add_slide(image_slide_layout)
        slide_corr.shapes.title.text = f"Correlation Heatmap: {task_name}"
        left = Inches(1)
        top = Inches(1.5)
        height = Inches(4)
        slide_corr.shapes.add_picture(fig_path, left, top, height=height)
    prs.save(output_pptx)
    print(f"PowerPoint saved to {output_pptx}")

from ecodynamics_multimodel_framework import EcodynamicsAI

framework = EcodynamicsAI()
framework.load_data()
framework.load_otu_data()
framework.categorize_features()
framework.filter_columns_of_interest()
framework.validate_and_assign_otu_data()

# Define tasks
tasks = {
    "Microbiota Richness (Root) with Plant OTUs": {
        "target": "richness_microbiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Microbiota Richness (Leaf) with Plant OTUs": {
        "target": "richness_microbiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Microbiota Richness (Root) with Plant Metrics": {
        "target": "richness_microbiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Microbiota Richness (Leaf) with Plant Metrics": {
        "target": "richness_microbiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Microbiota Shannon (Root) with Plant OTUs": {
        "target": "Shannon_microbiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Microbiota Shannon (Leaf) with Plant OTUs": {
        "target": "Shannon_microbiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Microbiota Shannon (Root) with Plant Metrics": {
        "target": "Shannon_microbiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Microbiota Shannon (Leaf) with Plant Metrics": {
        "target": "Shannon_microbiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Pathobiota Richness (Root) with Plant OTUs": {
        "target": "richness_pathobiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Pathobiota Richness (Leaf) with Plant OTUs": {
        "target": "richness_pathobiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Pathobiota Richness (Root) with Plant Metrics": {
        "target": "richness_pathobiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Pathobiota Richness (Leaf) with Plant Metrics": {
        "target": "richness_pathobiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Pathobiota Shannon (Root) with Plant OTUs": {
        "target": "Shannon_pathobiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Pathobiota Shannon (Leaf) with Plant OTUs": {
        "target": "Shannon_pathobiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_otus']
    },
    "Pathobiota Shannon (Root) with Plant Metrics": {
        "target": "Shannon_pathobiota_root",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    },
    "Pathobiota Shannon (Leaf) with Plant Metrics": {
        "target": "Shannon_pathobiota_leaf",
        "features": framework.feature_categories['environmental'] + framework.feature_categories['plant_metrics']
    }
}

lofo_results, correlation_figs, dataset_sizes, feature_importance_figs = perform_lofo_and_correlation_plots(
    framework, tasks,
    top_k_corr=15,
    do_significance_test=True,
    test_mode=False
)

create_powerpoint_report(
    lofo_results=lofo_results,
    correlation_figs=correlation_figs,
    dataset_sizes=dataset_sizes,
    feature_importance_figs=feature_importance_figs,
    output_pptx="final_analysis_report.pptx"
)
